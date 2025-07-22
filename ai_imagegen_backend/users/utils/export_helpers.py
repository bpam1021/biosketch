import io
import logging
import requests
import tempfile
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from weasyprint import HTML
from django.conf import settings
import os

# Setup logger
logger = logging.getLogger(__name__)

# Optional: dynamic base URL fallback
BASE_URL = getattr(settings, "BASE_URL", "https://api.biosketch.ai")


def render_pptx(presentation_obj, slides=None):
    logger.info(f"[render_pptx] Start rendering presentation {presentation_obj.id}")
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    slide_width = prs.slide_width.inches
    slide_height = prs.slide_height.inches

    max_title_height = 1.0
    max_image_height = 3.5
    max_desc_height = 2.0

    slides = slides or presentation_obj.slides.all().order_by("order")

    for slide_obj in slides:
        slide = prs.slides.add_slide(blank_slide_layout)

        # Title
        if slide_obj.title:
            try:
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.3), Inches(9), Inches(max_title_height)
                )
                frame = title_box.text_frame
                frame.text = slide_obj.title
                frame.paragraphs[0].font.size = Pt(28)
                frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            except Exception as e:
                logger.warning(f"[Title Error] {e}")

        # Image
        img_url = slide_obj.image_url
        if img_url and img_url.startswith("/"):
            img_url = BASE_URL.rstrip("/") + img_url

        if img_url:
            try:
                logger.info(f"[Image] Downloading from: {img_url}")
                response = requests.get(img_url, timeout=10)
                img_bytes = io.BytesIO(response.content)
                img = Image.open(img_bytes)
                width_px, height_px = img.size
                aspect_ratio = width_px / height_px

                max_width_in = 7.5
                max_height_in = max_image_height

                if aspect_ratio > 1:
                    width_in = min(max_width_in, slide_width - 2)
                    height_in = width_in / aspect_ratio
                else:
                    height_in = min(max_height_in, slide_height - 3)
                    width_in = height_in * aspect_ratio

                left = (slide_width - width_in) / 2

                slide.shapes.add_picture(
                    img_bytes,
                    Inches(left),
                    Inches(1.2),
                    width=Inches(width_in),
                    height=Inches(height_in),
                )
            except UnidentifiedImageError:
                logger.warning(f"[Image Error] Unidentified image: {img_url}")
            except Exception as e:
                logger.error(f"[Image Error] {e}")

        # Description
        if slide_obj.description:
            try:
                textbox = slide.shapes.add_textbox(
                    Inches(0.7), Inches(slide_height - max_desc_height - 0.3), Inches(8.6), Inches(max_desc_height)
                )
                frame = textbox.text_frame
                frame.word_wrap = True
                frame.text = slide_obj.description
                p = frame.paragraphs[0]
                p.font.size = Pt(16)
                p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
            except Exception as e:
                logger.warning(f"[Description Error] {e}")

    try:
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)

        # Cross-platform temp debug file
        tmp_dir = tempfile.gettempdir()
        debug_path = os.path.join(tmp_dir, f"debug_{presentation_obj.id}.pptx")
        with open(debug_path, "wb") as f:
            f.write(buffer.getvalue())
        logger.info(f"[render_pptx] Saved debug PPTX to: {debug_path}")

        return buffer
    except Exception as e:
        logger.error(f"[Final Save Error] {e}")
        raise


def render_pdf(presentation_obj, slides=None):
    logger.info(f"[render_pdf] Start rendering PDF for presentation {presentation_obj.id}")

    html = """
    <html>
    <head>
        <style>
            body { font-family: Arial, sans-serif; margin: 40px; }
            .slide {
                page-break-after: always;
                text-align: center;
                max-width: 800px;
                margin: 0 auto;
            }
            h2 {
                font-size: 24pt;
                margin-bottom: 20px;
                color: #333;
            }
            img {
                max-width: 90%;
                max-height: 400px;
                height: auto;
                margin-bottom: 20px;
                border: 1px solid #ccc;
                border-radius: 4px;
            }
            p {
                font-size: 14pt;
                color: #555;
                text-align: left;
                padding: 0 40px;
            }
        </style>
    </head>
    <body>
    """

    slides = slides or presentation_obj.slides.all().order_by("order")

    for slide in slides:
        html += "<div class='slide'>"
        html += f"<h2>{slide.title or ''}</h2>"
        img_url = slide.image_url
        if img_url and img_url.startswith("/"):
            img_url = BASE_URL.rstrip("/") + img_url
        if img_url:
            html += f"<img src='{img_url}'><br>"
        if slide.description:
            html += f"<p>{slide.description}</p>"
        html += "</div>"

    html += "</body></html>"

    try:
        pdf_stream = io.BytesIO()
        HTML(string=html).write_pdf(pdf_stream)
        pdf_stream.seek(0)

        # Cross-platform temp debug file
        tmp_dir = tempfile.gettempdir()
        debug_path = os.path.join(tmp_dir, f"debug_{presentation_obj.id}.pdf")
        with open(debug_path, "wb") as f:
            f.write(pdf_stream.getvalue())
        logger.info(f"[render_pdf] Saved debug PDF to: {debug_path}")

        return pdf_stream
    except Exception as e:
        logger.error(f"[render_pdf Error] {e}")
        raise
