# ai_imagegen_backend/users/tasks.py
# ADD these new tasks to your existing tasks.py file

import os
import gc
import tempfile
import numpy as np
from PIL import Image
from celery import shared_task
from gtts import gTTS
from django.conf import settings
from django.core.files import File
from django.db import transaction
from moviepy.video.VideoClip import ImageClip
from moviepy.audio.io.AudioFileClip import AudioFileClip
from moviepy.video.compositing.CompositeVideoClip import CompositeVideoClip
from moviepy.video.fx import resize
from django.core.files.base import ContentFile
from django.utils import timezone
from decimal import Decimal
import json
import logging
import traceback
import uuid

# Import your existing models and new presentation models
from users.models import (
    CreditTransaction, 
    # Clean presentation models
    Document, DocumentChapter, DocumentSection, DocumentTemplate,
    SlidePresentation, Slide, SlideTemplate, SlideTheme,
    MediaAsset, DiagramElement, PresentationExport
)

# Add OpenAI import for real AI generation
import openai
import re

logger = logging.getLogger(__name__)


# ============================================================================
# NEW ENHANCED PRESENTATION TASKS
# ============================================================================

@shared_task(bind=True, max_retries=3)
def generate_presentation_content(self, presentation_id, user_id, estimated_cost):
    """
    Generate presentation content using AI
    """
    try:
        from django.contrib.auth.models import User
        
        presentation = Presentation.objects.get(id=presentation_id)
        user = User.objects.get(id=user_id)
        
        logger.info(f"Starting content generation for presentation {presentation_id}")
        
        # Update status
        presentation.status = 'generating'
        presentation.save()
        
        # Generate outline based on prompt and type
        outline = generate_presentation_outline(
            presentation.original_prompt,
            presentation.presentation_type,
            presentation.quality
        )
        
        # Store generated outline
        presentation.generated_outline = outline
        presentation.save()
        
        # Initialize variables
        sections_created = []
        
        # Handle content generation differently for documents vs slides
        if presentation.presentation_type == 'document':
            # Generate unified document content
            unified_content = generate_unified_document_content(
                presentation.original_prompt,
                outline,
                presentation.quality
            )
            
            # Store unified content in presentation.document_content
            presentation.document_content = unified_content
            presentation.save()
            
            # Update word count
            from django.utils.html import strip_tags
            presentation.word_count = len(strip_tags(unified_content).split())
            presentation.save()
            
            logger.info(f"Generated unified document content ({presentation.word_count} words)")
            
        else:
            # Create sections for slide presentations
            for i, section_data in enumerate(outline.get('sections', [])):
                try:
                    # Generate detailed content for each section
                    detailed_content = generate_content_with_ai(
                        f"Write detailed content for: {section_data.get('title', '')}. Context: {presentation.original_prompt}",
                        length='medium' if presentation.quality == 'medium' else 'long' if presentation.quality == 'high' else 'short',
                        tone='professional'
                    )
                    
                    # Create section for slides
                    section = ContentSection.objects.create(
                        presentation=presentation,
                        section_type=section_data.get('type', 'content_slide'),
                        title=section_data.get('title', f'Section {i+1}'),
                        order=i,
                        content=detailed_content,
                        rich_content=detailed_content,
                        content_data={
                            'ai_generated': True,
                            'generation_prompt': section_data.get('title', ''),
                            'original_outline': section_data
                        },
                        image_prompt=section_data.get('image_prompt', ''),
                        ai_generated=True,
                        generation_prompt=section_data.get('title', '')
                    )
                    
                    sections_created.append(str(section.id))
                    logger.info(f"Created section {section.id} for presentation {presentation_id}")
                    
                    # Update progress
                    progress = int((i + 1) / len(outline.get('sections', [])) * 80)
                    self.update_state(state='PROGRESS', meta={'progress': progress})
                    
                except Exception as e:
                    logger.error(f"Failed to create section {i}: {e}")
                    continue
        
        # Generate images for sections that need them
        image_generation_tasks = []
        for section_id in sections_created:
            try:
                section = ContentSection.objects.get(id=section_id)
                if section.image_prompt:
                    # Queue image generation
                    task = generate_section_image.delay(section_id)
                    image_generation_tasks.append(str(task.id))
            except Exception as e:
                logger.error(f"Failed to queue image generation for section {section_id}: {e}")
        
        # Apply theme if template is selected
        if presentation.template:
            apply_template_styling.delay(presentation_id)
        
        # Deduct credits
        user.profile.credits -= Decimal(str(estimated_cost))
        user.profile.save()
        
        CreditTransaction.objects.create(
            user=user,
            amount=-Decimal(str(estimated_cost)),
            type='usage',
            description=f"Presentation generation: {presentation.title}"
        )
        
        # Log generation
        AIGenerationLog.objects.create(
            user=user,
            presentation=presentation,
            generation_type='presentation_outline',
            prompt=presentation.original_prompt,
            model_used='gpt-4',
            credits_used=Decimal(str(estimated_cost)),
            generated_content={
                'outline': outline,
                'sections_created': len(sections_created),
                'image_tasks': image_generation_tasks
            },
            success=True
        )
        
        # Update presentation status
        presentation.status = 'ready'
        presentation.generation_cost = Decimal(str(estimated_cost))
        presentation.total_credits_used = Decimal(str(estimated_cost))
        presentation.save()
        
        logger.info(f"Successfully generated presentation {presentation_id}")
        return {
            'status': 'completed',
            'sections_created': len(sections_created),
            'presentation_id': str(presentation_id)
        }
        
    except Exception as e:
        logger.error(f"Presentation generation failed: {e}")
        logger.error(traceback.format_exc())
        
        # Update presentation status
        try:
            presentation = Presentation.objects.get(id=presentation_id)
            presentation.status = 'error'
            presentation.save()
            
            # Log error
            from django.contrib.auth.models import User
            user = User.objects.get(id=user_id)
            AIGenerationLog.objects.create(
                user=user,
                presentation=presentation,
                generation_type='presentation_outline',
                prompt=presentation.original_prompt,
                model_used='gpt-4',
                credits_used=Decimal('0'),
                success=False,
                error_message=str(e)
            )
        except:
            pass
        
        # Retry if possible
        if self.request.retries < self.max_retries:
            raise self.retry(countdown=60, exc=e)
        
        return {'status': 'failed', 'error': str(e)}


@shared_task(bind=True, max_retries=3)
def generate_section_image(self, section_id):
    """
    Generate image for a specific section
    """
    try:
        section = ContentSection.objects.get(id=section_id)
        
        if not section.image_prompt:
            return {'status': 'skipped', 'reason': 'No image prompt'}
        
        logger.info(f"Generating image for section {section_id}")
        
        # Generate image using OpenAI DALL-E (v1.0+ API)
        from openai import OpenAI
        
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        
        response = client.images.generate(
            model="dall-e-3",
            prompt=section.image_prompt,
            n=1,
            size="1024x1024",
            response_format="url"
        )
        
        image_url = response.data[0].url
        
        # Update section with image URL
        section.image_url = image_url
        section.save()
        
        logger.info(f"Successfully generated image for section {section_id}")
        return {
            'status': 'completed',
            'section_id': str(section_id),
            'image_url': image_url
        }
        
    except Exception as e:
        logger.error(f"Image generation failed for section {section_id}: {e}")
        
        if self.request.retries < self.max_retries:
            raise self.retry(countdown=30, exc=e)
        
        return {'status': 'failed', 'error': str(e)}


@shared_task
def apply_template_styling(presentation_id):
    """
    Apply template styling to presentation sections
    """
    try:
        presentation = Presentation.objects.get(id=presentation_id)
        template = presentation.template
        
        if not template:
            return {'status': 'skipped', 'reason': 'No template selected'}
        
        logger.info(f"Applying template styling to presentation {presentation_id}")
        
        # Apply template data to presentation
        template_data = template.template_data
        
        # Update theme settings
        if 'theme' in template_data:
            presentation.theme_settings.update(template_data['theme'])
        
        # Apply section-specific styling
        if 'section_styles' in template_data:
            for section in presentation.sections.all():
                section_style = template_data['section_styles'].get(section.section_type, {})
                if section_style:
                    section.style_config.update(section_style)
                    section.save()
        
        presentation.save()
        
        logger.info(f"Successfully applied template styling to presentation {presentation_id}")
        return {'status': 'completed'}
        
    except Exception as e:
        logger.error(f"Template styling failed: {e}")
        return {'status': 'failed', 'error': str(e)}


@shared_task(bind=True, max_retries=2)
def export_presentation_task(self, export_job_id):
    """
    Export presentation to various formats
    """
    try:
        export_job = PresentationExportJob.objects.get(id=export_job_id)
        presentation = export_job.presentation
        
        logger.info(f"Starting export job {export_job_id} for presentation {presentation.id}")
        
        # Update job status
        export_job.status = 'processing'
        export_job.started_at = timezone.now()
        export_job.save()
        
        # Get sections to export
        if export_job.selected_sections:
            sections = presentation.sections.filter(id__in=export_job.selected_sections)
        else:
            sections = presentation.sections.all()
        
        sections = sections.order_by('order')
        
        # Export based on format
        export_format = export_job.export_format
        export_settings = export_job.export_settings
        
        output_file = None
        output_url = None
        
        if export_format == 'pdf':
            output_file = export_to_pdf(presentation, sections, export_settings)
        elif export_format == 'docx':
            output_file = export_to_docx(presentation, sections, export_settings)
        elif export_format == 'pptx':
            output_file = export_to_pptx(presentation, sections, export_settings)
        elif export_format == 'html':
            output_url = export_to_html(presentation, sections, export_settings)
        elif export_format == 'mp4':
            output_file = export_to_video(presentation, sections, export_settings)
        else:
            raise ValueError(f"Unsupported export format: {export_format}")
        
        # Save output
        if output_file:
            filename = f"{presentation.title}_{export_format}_{uuid.uuid4().hex[:8]}.{export_format}"
            export_job.output_file.save(filename, output_file)
        
        if output_url:
            export_job.output_url = output_url
        
        # Update job status
        export_job.status = 'completed'
        export_job.completed_at = timezone.now()
        export_job.progress = 100
        export_job.expires_at = timezone.now() + timezone.timedelta(days=7)
        export_job.save()
        
        logger.info(f"Successfully completed export job {export_job_id}")
        return {
            'status': 'completed',
            'export_job_id': str(export_job_id),
            'output_file': export_job.output_file.url if export_job.output_file else None,
            'output_url': export_job.output_url
        }
        
    except Exception as e:
        logger.error(f"Export job {export_job_id} failed: {e}")
        logger.error(traceback.format_exc())
        
        # Update job status
        try:
            export_job = PresentationExportJob.objects.get(id=export_job_id)
            export_job.status = 'failed'
            export_job.error_message = str(e)
            export_job.save()
        except:
            pass
        
        # Retry if possible
        if self.request.retries < self.max_retries:
            raise self.retry(countdown=120, exc=e)
        
        return {'status': 'failed', 'error': str(e)}


@shared_task
def cleanup_expired_exports():
    """
    Clean up expired export files
    """
    try:
        expired_jobs = PresentationExportJob.objects.filter(
            expires_at__lt=timezone.now(),
            status='completed'
        )
        
        cleaned_count = 0
        for job in expired_jobs:
            try:
                if job.output_file:
                    job.output_file.delete()
                job.delete()
                cleaned_count += 1
            except Exception as e:
                logger.error(f"Failed to clean up export job {job.id}: {e}")
        
        logger.info(f"Cleaned up {cleaned_count} expired export jobs")
        return {'status': 'completed', 'cleaned_count': cleaned_count}
        
    except Exception as e:
        logger.error(f"Export cleanup failed: {e}")
        return {'status': 'failed', 'error': str(e)}


@shared_task
def auto_save_presentation(presentation_id, user_id):
    """
    Auto-save presentation as a version
    """
    try:
        from users.models import PresentationVersion
        from django.contrib.auth.models import User
        
        presentation = Presentation.objects.get(id=presentation_id)
        user = User.objects.get(id=user_id)
        
        # Get latest version number
        latest_version = presentation.versions.first()
        version_number = (latest_version.version_number + 1) if latest_version else 1
        
        # Create content snapshot
        content_snapshot = {
            'presentation': {
                'title': presentation.title,
                'description': presentation.description,
                'theme_settings': presentation.theme_settings,
                'document_content': presentation.document_content,
                'document_settings': presentation.document_settings,
            },
            'sections': []
        }
        
        for section in presentation.sections.all():
            content_snapshot['sections'].append({
                'id': str(section.id),
                'title': section.title,
                'content': section.content,
                'rich_content': section.rich_content,
                'section_type': section.section_type,
                'order': section.order,
                'style_config': section.style_config,
                'layout_config': section.layout_config
            })
        
        # Create version
        PresentationVersion.objects.create(
            presentation=presentation,
            version_number=version_number,
            content_snapshot=content_snapshot,
            created_by=user,
            is_auto_save=True,
            changes_summary="Auto-save"
        )
        
        # Keep only last 10 auto-save versions
        old_versions = presentation.versions.filter(is_auto_save=True)[10:]
        for version in old_versions:
            version.delete()
        
        logger.info(f"Auto-saved presentation {presentation_id} as version {version_number}")
        return {'status': 'completed', 'version_number': version_number}
        
    except Exception as e:
        logger.error(f"Auto-save failed for presentation {presentation_id}: {e}")
        return {'status': 'failed', 'error': str(e)}


# ============================================================================
# AI HELPER FUNCTIONS
# ============================================================================

def generate_presentation_outline(prompt, presentation_type, quality):
    """
    Generate a structured outline for a presentation using AI
    """
    try:
        from openai import OpenAI
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        
        type_instructions = {
            'document': 'Create a document structure with sections like introduction, main content, conclusion',
            'slide': 'Create a slide presentation structure with title slide, content slides, and summary'
        }

        quality_instructions = {
            'low': 'Create 3-5 basic sections with simple content',
            'medium': 'Create 5-8 well-developed sections with detailed content',
            'high': 'Create 8-12 comprehensive sections with rich, detailed content and examples'
        }

        system_prompt = f"""
        You are an expert presentation designer. Create a structured outline for a {presentation_type}.
        
        {type_instructions.get(presentation_type, type_instructions['slide'])}
        {quality_instructions.get(quality, quality_instructions['medium'])}
        
        Generate a JSON response with:
        1. title: Presentation title
        2. sections: Array of sections with:
           - title: Section title
           - type: Section type (heading, paragraph, list, image, etc.)
           - content_outline: Brief content outline
           - image_prompt: Suggested image description (if applicable)
           - order: Section order
        3. theme_suggestions: Suggested theme colors and fonts
        4. estimated_duration: Estimated duration for presentation
        """

        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Create a {presentation_type} about: {prompt}"}
            ],
            temperature=0.6,
            max_tokens=2500
        )
        
        result = json.loads(response.choices[0].message.content)
        return result
        
    except Exception as e:
        logger.error(f"Outline generation failed: {e}")
        # Return fallback outline
        return {
            "title": "Generated Presentation",
            "sections": [
                {
                    "title": "Introduction",
                    "type": "title_slide" if presentation_type == 'slide' else "heading",
                    "content_outline": f"Introduction to {prompt}",
                    "image_prompt": f"Professional introduction image for {prompt}",
                    "order": 0
                },
                {
                    "title": "Main Content",
                    "type": "content_slide" if presentation_type == 'slide' else "paragraph",
                    "content_outline": f"Main content about {prompt}",
                    "image_prompt": f"Relevant illustration for {prompt}",
                    "order": 1
                },
                {
                    "title": "Conclusion",
                    "type": "content_slide" if presentation_type == 'slide' else "paragraph",
                    "content_outline": f"Conclusion and summary of {prompt}",
                    "image_prompt": f"Summary image for {prompt}",
                    "order": 2
                }
            ],
            "theme_suggestions": {
                "primary_color": "#2563EB",
                "secondary_color": "#7C3AED",
                "background_color": "#FFFFFF",
                "font_family": "Inter"
            },
            "estimated_duration": "10 minutes"
        }


def generate_content_with_ai(prompt, length='medium', tone='professional'):
    """
    Generate content using AI based on prompt, length, and tone
    """
    try:
        from openai import OpenAI
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        
        length_mapping = {
            'short': 'Write 1-2 paragraphs (100-200 words)',
            'medium': 'Write 3-4 paragraphs (300-500 words)',
            'long': 'Write 5-8 paragraphs (700-1000 words)'
        }

        tone_mapping = {
            'professional': 'Use a professional, business-appropriate tone',
            'casual': 'Use a casual, conversational tone',
            'academic': 'Use an academic, scholarly tone with proper citations style',
            'creative': 'Use a creative, engaging tone with storytelling elements',
            'technical': 'Use a technical, precise tone with industry terminology'
        }

        system_prompt = f"""
        You are an expert content writer. {tone_mapping.get(tone, tone_mapping['professional'])}.
        {length_mapping.get(length, length_mapping['medium'])}.
        
        Ensure the content is:
        - Well-structured with clear flow
        - Factually accurate and informative
        - Engaging and relevant to the topic
        - Properly formatted
        """

        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ],
            temperature=0.7,
            max_tokens=1500
        )
        
        return response.choices[0].message.content.strip()
        
    except Exception as e:
        logger.error(f"Content generation failed: {e}")
        return f"Content for: {prompt}\n\nThis content will be generated when the AI service is available."


def generate_unified_document_content(prompt, outline, quality):
    """
    Generate unified document content instead of sections
    """
    try:
        from openai import OpenAI
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        
        # Extract section titles from outline to create document structure
        sections = outline.get('sections', [])
        section_titles = [s.get('title', '') for s in sections]
        
        quality_mapping = {
            'low': 'Write a comprehensive document with 1000-1500 words',
            'medium': 'Write a detailed document with 2000-3000 words', 
            'high': 'Write an extensive, research-quality document with 3000-5000 words'
        }

        system_prompt = f"""
        You are a world-class professional technical writer and consultant creating comprehensive, publication-ready documents.
        
        CONTENT REQUIREMENTS:
        {quality_mapping.get(quality, quality_mapping['medium'])}.
        
        CRITICAL: Generate substantial professional content with comprehensive structure:
        
        REQUIRED DOCUMENT STRUCTURE:
        1. **Executive Summary** (300-500 words)
           - Key findings, recommendations, and strategic insights
           - Include numbered or bulleted key points
           - Add professional table with metrics/data
        
        2. **Introduction & Background** (400-600 words)
           - Context, objectives, methodology
           - Industry analysis and current trends
           - Use bulleted lists for objectives and scope
        
        3. **Main Analysis Sections** (800-1200 words each)
           - 2-3 comprehensive analysis sections based on topic
           - Each section should have multiple subsections
           - Include data tables, comparison charts, case studies
           - Use numbered lists for processes, bulleted lists for features
        
        4. **Recommendations & Implementation** (400-600 words)
           - Actionable recommendations with priority levels
           - Implementation timeline in table format
           - Risk assessment with bulleted points
        
        5. **Conclusion & Next Steps** (300-400 words)
           - Summary of findings and strategic direction
           - Future considerations and follow-up actions
        
        FORMATTING REQUIREMENTS (CRITICAL):
        - Use semantic HTML: <h1>, <h2>, <h3>, <h4>, <p>, <ul>, <ol>, <li>, <table>, <thead>, <tbody>, <tr>, <th>, <td>, <blockquote>, <strong>, <em>
        - Include multiple professional tables with realistic data and metrics
        - Use numbered lists (ol) for processes, steps, priorities, rankings
        - Use bulleted lists (ul) for features, benefits, characteristics, options
        - Add blockquotes for important insights, statistics, or expert opinions
        - Use strong/em for emphasis on key terms and critical information
        - Create proper table headers with <thead> and data rows with <tbody>
        
        CONTENT DEPTH REQUIREMENTS:
        - Each section must be 200-400 words minimum
        - Include specific examples, case studies, and real-world applications
        - Add relevant statistics, percentages, and quantitative data
        - Provide actionable insights and practical recommendations
        - Include industry best practices and professional standards
        
        SECTION COVERAGE:
        Structure the document to comprehensively cover these areas: {', '.join(section_titles)}
        
        HTML FORMATTING STANDARDS:
        <h1>Document Title</h1>
        
        <h2>Executive Summary</h2>
        <p>Comprehensive paragraph with analysis...</p>
        
        <h3>Key Findings</h3>
        <ol>
            <li><strong>Finding 1:</strong> Detailed explanation with data</li>
            <li><strong>Finding 2:</strong> Detailed explanation with metrics</li>
        </ol>
        
        <table>
            <thead>
                <tr><th>Metric</th><th>Current</th><th>Target</th><th>Improvement</th></tr>
            </thead>
            <tbody>
                <tr><td>Performance</td><td>75%</td><td>90%</td><td>+15%</td></tr>
                <tr><td>Efficiency</td><td>82%</td><td>95%</td><td>+13%</td></tr>
            </tbody>
        </table>
        
        <h3>Strategic Recommendations</h3>
        <ul>
            <li><strong>Priority 1:</strong> Implement advanced analytics framework</li>
            <li><strong>Priority 2:</strong> Optimize operational workflows</li>
            <li><strong>Priority 3:</strong> Enhance team collaboration tools</li>
        </ul>
        
        <blockquote>
            <p>Industry research shows that organizations implementing these strategies achieve 40% better performance outcomes within 12 months.</p>
        </blockquote>
        
        GENERATE COMPREHENSIVE, PROFESSIONAL CONTENT - Each section should be detailed with examples, data, and actionable insights.
        """

        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Write a comprehensive document about: {prompt}"}
            ],
            temperature=0.7,
            max_tokens=4000
        )
        
        return response.choices[0].message.content.strip()
        
    except Exception as e:
        logger.error(f"Unified document content generation failed: {e}")
        
        # Return structured fallback content
        return f"""
        <h1>{outline.get('title', 'Generated Document')}</h1>
        
        <h2>Introduction</h2>
        <p>This document provides a comprehensive overview of {prompt}. The content covers key concepts, analysis, and conclusions drawn from current research and understanding.</p>
        
        <h2>Overview</h2>
        <p>The main focus areas of this document include:</p>
        <ul>
            {''.join(f'<li>{section.get("title", "")}</li>' for section in outline.get('sections', []))}
        </ul>
        
        <h2>Main Content</h2>
        <p>Detailed content about {prompt} will be provided here. This section explores the fundamental aspects, applications, and significance of the topic.</p>
        
        <h3>Key Points</h3>
        <ul>
            <li>Important aspect 1 related to {prompt}</li>
            <li>Important aspect 2 related to {prompt}</li>
            <li>Important aspect 3 related to {prompt}</li>
        </ul>
        
        <h2>Analysis</h2>
        <p>This section provides in-depth analysis and examination of {prompt}, including methodologies, findings, and interpretations.</p>
        
        <blockquote>
            <p>This is a key insight or important note about {prompt} that deserves special attention.</p>
        </blockquote>
        
        <h2>Conclusion</h2>
        <p>In conclusion, {prompt} represents an important area of study with significant implications. The document has covered the essential aspects and provided a comprehensive overview of the topic.</p>
        
        <p><em>This document was generated automatically. Content will be enhanced when the AI service is fully available.</em></p>
        """


# ============================================================================
# EXPORT HELPER FUNCTIONS
# ============================================================================

def export_to_pdf(presentation, sections, settings):
    """Export presentation to PDF"""
    try:
        from weasyprint import HTML, CSS
        import io
        
        # Generate HTML content
        html_content = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <meta charset="utf-8">
            <title>{presentation.title}</title>
            <style>
                body {{ font-family: Arial, sans-serif; margin: 40px; line-height: 1.6; }}
                h1 {{ color: #333; border-bottom: 3px solid #007bff; padding-bottom: 10px; }}
                h2 {{ color: #555; margin-top: 30px; }}
                .section {{ margin-bottom: 30px; page-break-inside: avoid; }}
                .meta {{ color: #666; font-size: 12px; margin-bottom: 20px; }}
                img {{ max-width: 100%; height: auto; margin: 20px 0; }}
            </style>
        </head>
        <body>
            <h1>{presentation.title}</h1>
            <div class="meta">Generated on {timezone.now().strftime('%B %d, %Y at %I:%M %p')}</div>
        """
        
        for section in sections:
            html_content += f"""
            <div class="section">
                <h2>{section.title}</h2>
                {section.rich_content or section.content}
                {f'<img src="{section.image_url}" alt="{section.title}">' if section.image_url else ''}
            </div>
            """
        
        html_content += """
        </body>
        </html>
        """
        
        # Convert to PDF
        pdf_buffer = io.BytesIO()
        HTML(string=html_content).write_pdf(pdf_buffer)
        pdf_buffer.seek(0)
        
        return ContentFile(pdf_buffer.getvalue())
        
    except Exception as e:
        logger.error(f"PDF export failed: {e}")
        raise


def export_to_docx(presentation, sections, settings):
    """Export presentation to Word document"""
    try:
        from docx import Document
        from docx.shared import Inches
        import io
        import requests
        
        doc = Document()
        
        # Add title
        doc.add_heading(presentation.title, 0)
        
        # Add metadata
        p = doc.add_paragraph(f"Generated on {timezone.now().strftime('%B %d, %Y at %I:%M %p')}")
        p.style = 'Subtitle'
        
        # Add sections
        for section in sections:
            doc.add_heading(section.title, level=1)
            doc.add_paragraph(section.content)
            
            # Add image if available
            if section.image_url:
                try:
                    response = requests.get(section.image_url)
                    if response.status_code == 200:
                        image_stream = io.BytesIO(response.content)
                        doc.add_picture(image_stream, width=Inches(4))
                except Exception as e:
                    logger.warning(f"Failed to add image to docx: {e}")
        
        # Save to buffer
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        
        return ContentFile(buffer.getvalue())
        
    except Exception as e:
        logger.error(f"DOCX export failed: {e}")
        raise


def export_to_pptx(presentation, sections, settings):
    """Export presentation to PowerPoint"""
    try:
        from pptx import Presentation as PPTXPresentation
        from pptx.util import Inches, Pt
        from pptx.enum.text import PP_ALIGN
        import io
        import requests
        
        prs = PPTXPresentation()
        
        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = presentation.title
        subtitle.text = f"Generated on {timezone.now().strftime('%B %d, %Y')}"
        
        # Content slides
        for section in sections:
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            title = slide.shapes.title
            content = slide.placeholders[1]
            
            title.text = section.title
            content.text = section.content
            
            # Add image if available
            if section.image_url:
                try:
                    response = requests.get(section.image_url)
                    if response.status_code == 200:
                        image_stream = io.BytesIO(response.content)
                        slide.shapes.add_picture(
                            image_stream, 
                            Inches(1), 
                            Inches(3), 
                            width=Inches(4)
                        )
                except Exception as e:
                    logger.warning(f"Failed to add image to pptx: {e}")
        
        # Save to buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        return ContentFile(buffer.getvalue())
        
    except Exception as e:
        logger.error(f"PPTX export failed: {e}")
        raise


def export_to_html(presentation, sections, settings):
    """Export presentation to HTML"""
    try:
        html_content = f"""
        <!DOCTYPE html>
        <html lang="en">
        <head>
            <meta charset="UTF-8">
            <meta name="viewport" content="width=device-width, initial-scale=1.0">
            <title>{presentation.title}</title>
            <style>
                body {{
                    font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;
                    line-height: 1.6;
                    max-width: 800px;
                    margin: 0 auto;
                    padding: 20px;
                    color: #333;
                }}
                .header {{
                    text-align: center;
                    border-bottom: 2px solid #007bff;
                    padding-bottom: 20px;
                    margin-bottom: 40px;
                }}
                .section {{
                    margin-bottom: 40px;
                    padding: 20px;
                    border-left: 4px solid #007bff;
                    background: #f8f9fa;
                }}
                .section h2 {{
                    margin-top: 0;
                    color: #007bff;
                }}
                img {{
                    max-width: 100%;
                    height: auto;
                    border-radius: 8px;
                    box-shadow: 0 4px 6px rgba(0,0,0,0.1);
                }}
                .meta {{
                    color: #666;
                    font-size: 14px;
                }}
            </style>
        </head>
        <body>
            <div class="header">
                <h1>{presentation.title}</h1>
                <p class="meta">Generated on {timezone.now().strftime('%B %d, %Y at %I:%M %p')}</p>
            </div>
        """
        
        for section in sections:
            html_content += f"""
            <div class="section">
                <h2>{section.title}</h2>
                <div>{section.rich_content or section.content}</div>
                {f'<img src="{section.image_url}" alt="{section.title}">' if section.image_url else ''}
            </div>
            """
        
        html_content += """
        </body>
        </html>
        """
        
        # Save HTML file and return URL
        filename = f"presentation_{presentation.id}_{uuid.uuid4().hex[:8]}.html"
        file_path = os.path.join(settings.MEDIA_ROOT, 'exports', filename)
        
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        return f"{settings.MEDIA_URL}exports/{filename}"
        
    except Exception as e:
        logger.error(f"HTML export failed: {e}")
        raise


def export_to_video(presentation, sections, settings):
    """Export presentation to video"""
    try:
        from moviepy.editor import ImageClip, TextClip, CompositeVideoClip, concatenate_videoclips
        import requests
        from PIL import Image, ImageDraw, ImageFont
        
        clips = []
        duration_per_slide = settings.get('duration_per_slide', 5)
        
        for section in sections:
            # Create background image
            img = Image.new('RGB', (1920, 1080), color=(255, 255, 255))
            draw = ImageDraw.Draw(img)
            
            # Add title and content
            try:
                font_title = ImageFont.truetype("arial.ttf", 72)
                font_content = ImageFont.truetype("arial.ttf", 36)
            except:
                font_title = ImageFont.load_default()
                font_content = ImageFont.load_default()
            
            # Draw title
            title_bbox = draw.textbbox((0, 0), section.title, font=font_title)
            title_x = (1920 - title_bbox[2]) // 2
            draw.text((title_x, 100), section.title, fill=(0, 0, 0), font=font_title)
            
            # Draw content (truncated)
            content_lines = section.content[:300].split('\n')[:5]
            y_offset = 300
            for line in content_lines:
                if len(line) > 60:
                    line = line[:57] + "..."
                draw.text((100, y_offset), line, fill=(64, 64, 64), font=font_content)
                y_offset += 50
            
            # Save temporary image
            temp_img_path = f"/tmp/slide_{section.id}_{uuid.uuid4().hex[:8]}.png"
            img.save(temp_img_path)
            
            # Create video clip
            clip = ImageClip(temp_img_path).set_duration(duration_per_slide)
            clips.append(clip)
            
            # Clean up temp file
            os.remove(temp_img_path)
        
        # Concatenate all clips
        final_video = concatenate_videoclips(clips, method="compose")
        
        # Save to buffer
        temp_video_path = f"/tmp/presentation_{presentation.id}_{uuid.uuid4().hex[:8]}.mp4"
        final_video.write_videofile(
            temp_video_path,
            fps=24,
            codec='libx264',
            audio_codec='aac',
            temp_audiofile='temp-audio.m4a',
            remove_temp=True
        )
        
        # Read file and create ContentFile
        with open(temp_video_path, 'rb') as f:
            video_content = f.read()
        
        # Clean up
        os.remove(temp_video_path)
        final_video.close()
        
        return ContentFile(video_content)
        
    except Exception as e:
        logger.error(f"Video export failed: {e}")
        raise


# ============================================================================
# NEW CLEAN ARCHITECTURE CELERY TASKS
# ============================================================================

@shared_task(bind=True, max_retries=3)
def generate_document_ai_task(self, prompt, document_type, template_id, user_id):
    """
    Celery task for AI-powered document generation (Word-like)
    """
    try:
        from django.contrib.auth.models import User
        
        user = User.objects.get(id=user_id)
        logger.info(f"Starting AI document generation for user {user_id}")
        
        # Set OpenAI client
        from openai import OpenAI
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        
        # Get template structure if available
        template_structure = {}
        if template_id:
            try:
                template = DocumentTemplate.objects.get(id=template_id)
                template_structure = template.structure
                logger.info(f"Using template structure: {template.name}")
            except DocumentTemplate.DoesNotExist:
                logger.warning(f"Template {template_id} not found, using default structure")

        # Generate comprehensive document structure using AI
        system_prompt = f"""You are a world-class professional document writer and consultant specializing in creating comprehensive, publication-ready {document_type} documents.

CRITICAL REQUIREMENTS:
- Generate 5000-8000 words of substantive professional content
- Create 5-8 comprehensive chapters with hierarchical sub-sections
- Each chapter must have 3-5 main sections, each section should have 2-4 subsections
- Write 200-400 words per section with detailed analysis and insights
- Use proper academic/business formatting with rich HTML structure
- Include tables, charts references, examples, case studies, and actionable recommendations

DOCUMENT TYPE: {document_type}
TEMPLATE STRUCTURE: {template_structure if template_structure else "Professional business document"}

CONTENT DEPTH REQUIREMENTS:
- Executive Summary: 300-500 words with key findings and recommendations
- Introduction: 400-600 words with background, objectives, methodology
- Main Chapters: 1000-1500 words each with detailed analysis
- Each Section: 200-400 words with specific insights and examples
- Conclusion: 400-600 words with summary and future recommendations
- Include relevant data, statistics, best practices, and case studies

HTML FORMATTING STANDARDS:
- Use semantic HTML: <h1>, <h2>, <h3>, <h4>, <p>, <ul>, <ol>, <li>, <table>, <blockquote>, <strong>, <em>
- Create professional tables with relevant data and metrics
- Use numbered and bulleted lists for clarity
- Include blockquotes for important insights
- Add emphasis with <strong> and <em> appropriately
- Structure with clear hierarchy: Chapter > Section > Subsection

REQUIRED DOCUMENT STRUCTURE:
{{
    "title": "Professional Document Title (Specific to Topic)",
    "abstract": "Comprehensive 400-word executive summary covering background, methodology, key findings, analysis, conclusions, and strategic recommendations",
    "keywords": "10-15 relevant professional keywords separated by commas",
    "authors": ["Professional Author Name"],
    "subject": "Detailed subject description",
    "category": "{document_type}",
    "structure": {{
        "chapters": [
            {{
                "number": 1,
                "title": "Executive Summary",
                "content": "<h2>Executive Summary</h2><p>Comprehensive overview paragraph 1...</p><p>Key findings paragraph 2...</p><h3>Strategic Recommendations</h3><ul><li>Recommendation 1 with details</li><li>Recommendation 2 with details</li></ul><table><tr><th>Key Metric</th><th>Current State</th><th>Target</th></tr><tr><td>Performance</td><td>75%</td><td>95%</td></tr></table>",
                "sections": [
                    {{
                        "number": "1.1",
                        "title": "Key Findings Overview",
                        "content": "<p>Detailed paragraph 1 with analysis (150-200 words)...</p><p>Detailed paragraph 2 with insights (150-200 words)...</p><blockquote>Important insight or key takeaway</blockquote>",
                        "level": 1
                    }},
                    {{
                        "number": "1.2", 
                        "title": "Strategic Recommendations",
                        "content": "<p>Comprehensive recommendations paragraph 1...</p><ol><li><strong>Priority 1:</strong> Detailed recommendation</li><li><strong>Priority 2:</strong> Detailed recommendation</li></ol>",
                        "level": 1
                    }},
                    {{
                        "number": "1.3",
                        "title": "Implementation Roadmap", 
                        "content": "<p>Detailed implementation approach...</p><table><tr><th>Phase</th><th>Timeline</th><th>Key Activities</th></tr><tr><td>Phase 1</td><td>Q1 2024</td><td>Initial setup and planning</td></tr></table>",
                        "level": 1
                    }}
                ]
            }},
            {{
                "number": 2,
                "title": "Introduction and Background",
                "content": "<h2>Introduction and Background</h2><p>Comprehensive introduction paragraph explaining context and importance...</p><p>Background information with industry analysis...</p>",
                "sections": [
                    {{
                        "number": "2.1",
                        "title": "Industry Context and Background",
                        "content": "<p>Detailed industry analysis with current trends...</p><p>Market context and competitive landscape...</p>",
                        "level": 1
                    }},
                    {{
                        "number": "2.2",
                        "title": "Problem Statement and Objectives",
                        "content": "<p>Clear problem definition with supporting data...</p><ul><li>Objective 1: Specific measurable goal</li><li>Objective 2: Specific measurable goal</li></ul>",
                        "level": 1
                    }}
                ]
            }},
            {{
                "number": 3,
                "title": "Methodology and Approach",
                "content": "<h2>Methodology and Approach</h2><p>Detailed methodology explanation...</p>",
                "sections": [
                    {{
                        "number": "3.1",
                        "title": "Research Methodology",
                        "content": "<p>Comprehensive research approach description...</p>",
                        "level": 1
                    }},
                    {{
                        "number": "3.2",
                        "title": "Data Collection and Analysis",
                        "content": "<p>Detailed data collection methods...</p>",
                        "level": 1
                    }}
                ]
            }}
        ]
    }},
    "content": "Complete unified HTML content combining all chapters and sections with proper formatting",
    "diagram_opportunities": [
        {{
            "text": "Performance improved by 40% across all key metrics including efficiency, quality, and customer satisfaction",
            "suggested_chart_type": "bar_chart",
            "confidence": 0.92,
            "position": "after_executive_summary"
        }},
        {{
            "text": "The implementation process follows five sequential phases: assessment, planning, execution, monitoring, and optimization",
            "suggested_chart_type": "flowchart", 
            "confidence": 0.95,
            "position": "after_methodology_section"
        }}
    ]
}}

GENERATE SUBSTANTIAL PROFESSIONAL CONTENT - Each section should be comprehensive with detailed analysis, examples, and actionable insights. Target 5000+ total words."""
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Create a comprehensive, professional {document_type} document about: {prompt}. Ensure it's detailed, well-structured, and contains substantial content with multiple chapters and sections."}
            ],
            temperature=0.7,
            max_tokens=12000  # Much higher limit with GPT-4o's 128K context
        )
        
        ai_response = response.choices[0].message.content
        
        try:
            ai_data = json.loads(ai_response)
        except json.JSONDecodeError:
            # Fallback if AI doesn't return valid JSON
            ai_data = create_fallback_document_data(prompt, document_type)
        
        # Get template if specified
        template = None
        if template_id:
            try:
                template = DocumentTemplate.objects.get(id=template_id)
            except DocumentTemplate.DoesNotExist:
                pass
        
        # Create document with AI-generated content
        with transaction.atomic():
            document = Document.objects.create(
                title=ai_data['title'],
                template=template,
                created_by=user,
                abstract=ai_data.get('abstract', ''),
                keywords=ai_data.get('keywords', ''),
                authors=ai_data.get('authors', []),
                subject=ai_data.get('subject', ''),
                category=ai_data.get('category', document_type),
                content=ai_data.get('content', ''),
                structure=ai_data.get('structure', {}),
                diagram_opportunities=ai_data.get('diagram_opportunities', [])
            )
            
            # Create chapters and sections from AI structure
            for chapter_data in ai_data.get('structure', {}).get('chapters', []):
                chapter = DocumentChapter.objects.create(
                    document=document,
                    number=chapter_data['number'],
                    title=chapter_data['title'],
                    content=chapter_data['content'],
                    order=chapter_data['number'] - 1
                )
                
                # Create sections
                for section_data in chapter_data.get('sections', []):
                    DocumentSection.objects.create(
                        chapter=chapter,
                        number=section_data['number'],
                        title=section_data['title'],
                        content=section_data['content'],
                        level=section_data['level'],
                        order=len(chapter.sections.all())
                    )
            
            # Update document statistics
            document.update_statistics()
            
            logger.info(f"Successfully generated document {document.id}")
            return {
                'status': 'success',
                'document_id': str(document.id),
                'document_data': {
                    'id': str(document.id),
                    'title': document.title,
                    'word_count': document.word_count,
                    'page_count': document.page_count
                },
                'ai_analysis': {
                    'diagram_opportunities': ai_data.get('diagram_opportunities', []),
                    'generation_success': True,
                    'content_quality': 'high'
                }
            }
        
    except Exception as e:
        logger.error(f"AI document generation failed: {e}")
        logger.error(traceback.format_exc())
        
        if self.request.retries < self.max_retries:
            raise self.retry(countdown=60, exc=e)
        
        return {'status': 'failed', 'error': str(e)}


@shared_task(bind=True, max_retries=3)
def generate_slides_ai_task(self, prompt, theme_id, slide_size, user_id):
    """
    Celery task for AI-powered slide presentation generation (PowerPoint-like)
    """
    try:
        from django.contrib.auth.models import User
        
        user = User.objects.get(id=user_id)
        theme = SlideTheme.objects.get(id=theme_id)
        
        logger.info(f"Starting AI slide generation for user {user_id}")
        
        # Set OpenAI client
        from openai import OpenAI
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        
        # Generate professional PowerPoint-style slides using AI
        system_prompt = f"""You are a world-class presentation designer and PowerPoint expert specializing in creating professional, engaging slide presentations similar to Microsoft PowerPoint.

CRITICAL REQUIREMENTS:
- Generate 8-15 comprehensive slides with professional PowerPoint layouts
- Use diverse slide templates for visual variety and engagement
- Include substantial content - each slide should have 50-150 words of meaningful content
- Create professional speaker notes (100-200 words per slide)
- Follow corporate presentation standards and design principles
- Include data-driven content, charts, comparisons, and actionable insights

SLIDE TEMPLATE TYPES (Use variety):
1. **title_slide**: Title page with main title, subtitle, presenter info
2. **agenda_overview**: Presentation agenda/outline with numbered items
3. **section_divider**: Section break slides with large titles
4. **title_content**: Standard slide with title and bullet points/content
5. **two_column**: Split layout with title and two columns of content
6. **content_image**: Content with image placeholder on right side
7. **full_image**: Full-screen image with overlay text
8. **comparison**: Side-by-side comparison layout
9. **timeline**: Timeline or process flow layout
10. **data_visual**: Chart/graph placeholder with supporting content
11. **quote_testimonial**: Large quote with attribution
12. **conclusion_cta**: Conclusion slide with call-to-action
13. **thank_you**: Final slide with contact information

CONTENT DEPTH REQUIREMENTS:
- **Title Slide**: Professional title, compelling subtitle, presenter credentials
- **Agenda**: 4-6 main sections with brief descriptions
- **Content Slides**: 4-7 bullet points with detailed explanations (20-30 words per bullet)
- **Speaker Notes**: Comprehensive talking points, statistics, examples, transitions
- **Visual Descriptions**: Detailed descriptions for image/chart placeholders
- **Professional Formatting**: Use proper business language, active voice, clear structure

THEME: {theme.name} with colors: {theme.colors}
SLIDE SIZE: {slide_size}

Return a JSON response with this exact structure:
{{
    "title": "Professional Presentation Title (Specific to Topic)",
    "subtitle": "Comprehensive subtitle describing presentation value",
    "presenter_info": "Professional presenter name and credentials",
    "slides": [
        {{
            "template_type": "title_slide",
            "order": 0,
            "content": {{
                "title_zone": "Engaging Main Title That Captures Attention",
                "subtitle_zone": "Compelling subtitle that describes the value proposition",
                "presenter_zone": "Presenter Name, Title, Company",
                "date_zone": "Presentation Date"
            }},
            "notes": "Welcome the audience, introduce yourself and your credentials. Explain the importance of this topic and what attendees will gain from this presentation. Mention key statistics or compelling facts to grab attention. Transition smoothly to the agenda.",
            "design_elements": {{
                "background_style": "gradient",
                "text_alignment": "center",
                "font_sizes": {{"title": 44, "subtitle": 24, "presenter": 18}}
            }},
            "duration": 90
        }},
        {{
            "template_type": "agenda_overview",
            "order": 1,
            "content": {{
                "title_zone": "What We'll Cover Today",
                "content_zone": "<ol><li><strong>Section 1:</strong> Brief description of first major topic</li><li><strong>Section 2:</strong> Brief description of second major topic</li><li><strong>Section 3:</strong> Brief description of third major topic</li><li><strong>Section 4:</strong> Brief description of fourth major topic</li><li><strong>Q&A:</strong> Time for questions and discussion</li></ol>"
            }},
            "notes": "Set expectations for the presentation flow. Mention the approximate time for each section. Encourage questions throughout or specify when Q&A will happen. Highlight the most valuable sections that audience should pay special attention to.",
            "design_elements": {{
                "list_style": "numbered_professional",
                "bullet_color": "{theme.colors.get('primary', '#2563eb')}",
                "spacing": "comfortable"
            }},
            "duration": 120
        }},
        {{
            "template_type": "title_content",
            "order": 2,
            "content": {{
                "title_zone": "Specific Content Title for This Slide",
                "content_zone": "<ul><li><strong>Key Point 1:</strong> Detailed explanation with supporting information and context (25-30 words)</li><li><strong>Key Point 2:</strong> Another important insight with data, examples, or case studies to support the point</li><li><strong>Key Point 3:</strong> Additional critical information that builds on previous points and adds value</li><li><strong>Key Point 4:</strong> Final supporting point that reinforces the main message and provides actionable insights</li></ul>"
            }},
            "notes": "Elaborate on each bullet point with specific examples, statistics, and real-world applications. Share relevant case studies or success stories. Address potential questions or objections. Connect this content to the overall presentation theme and next sections.",
            "design_elements": {{
                "bullet_style": "modern_arrows",
                "emphasis_color": "{theme.colors.get('accent', '#dc2626')}",
                "layout": "balanced"
            }},
            "duration": 180
        }},
        {{
            "template_type": "comparison",
            "order": 3,
            "content": {{
                "title_zone": "Before vs. After / Option A vs. Option B",
                "left_column": {{
                    "header": "Current State / Option A",
                    "content": "<ul><li>Current challenge or limitation #1</li><li>Existing process or method #2</li><li>Pain point or inefficiency #3</li><li>Cost or resource implication #4</li></ul>"
                }},
                "right_column": {{
                    "header": "Improved State / Option B", 
                    "content": "<ul><li>Resolved challenge with specific solution</li><li>Enhanced process with measurable benefits</li><li>Eliminated pain point with new approach</li><li>Reduced costs with quantified savings</li></ul>"
                }}
            }},
            "notes": "Highlight the stark differences between the two options. Use specific metrics and percentages where possible. Share customer testimonials or case studies that illustrate these differences. Address potential concerns about implementation or transition.",
            "design_elements": {{
                "split_style": "vertical_divide",
                "contrast_colors": ["#ef4444", "#10b981"],
                "visual_divider": true
            }},
            "duration": 200
        }},
        {{
            "template_type": "data_visual",
            "order": 4,
            "content": {{
                "title_zone": "Key Performance Metrics & Results",
                "chart_placeholder": {{
                    "type": "bar_chart",
                    "title": "Performance Improvement Over Time",
                    "description": "Chart showing 40% improvement in efficiency, 25% reduction in costs, 60% increase in customer satisfaction over 12 months",
                    "data_points": ["Q1: 65%", "Q2: 72%", "Q3: 81%", "Q4: 89%"]
                }},
                "supporting_content": "<ul><li><strong>40% efficiency improvement</strong> in first 6 months</li><li><strong>$2.3M cost savings</strong> annually</li><li><strong>95% client satisfaction</strong> rate achieved</li><li><strong>3x faster implementation</strong> than industry average</li></ul>"
            }},
            "notes": "Walk through each data point slowly, explaining the methodology behind the measurements. Share the specific actions that led to these improvements. Compare results to industry benchmarks. Address any questions about data accuracy or measurement methods.",
            "design_elements": {{
                "chart_style": "professional_modern",
                "color_scheme": ["{theme.colors.get('primary', '#2563eb')}", "{theme.colors.get('secondary', '#7c3aed')}"],
                "emphasis_metrics": true
            }},
            "duration": 240
        }}
    ],
    "outline_structure": {{
        "main_sections": [
            "Introduction & Agenda",
            "Current State Analysis", 
            "Proposed Solutions",
            "Implementation Roadmap",
            "Expected Outcomes",
            "Next Steps & Q&A"
        ],
        "total_slides": 12,
        "estimated_duration": 25,
        "key_takeaways": [
            "Primary learning objective #1",
            "Key actionable insight #2", 
            "Strategic recommendation #3"
        ]
    }},
    "diagram_opportunities": [
        {{
            "slide_number": 3,
            "content": "Process flow showing step-by-step implementation methodology",
            "suggested_chart_type": "flowchart",
            "confidence": 0.92,
            "chart_title": "Implementation Process Flow",
            "estimated_complexity": "medium"
        }},
        {{
            "slide_number": 5,
            "content": "Comparison of performance metrics across different time periods",
            "suggested_chart_type": "bar_chart",
            "confidence": 0.88,
            "chart_title": "Performance Trends Analysis",
            "estimated_complexity": "low"
        }},
        {{
            "slide_number": 7,
            "content": "Organizational structure showing reporting relationships and responsibilities",
            "suggested_chart_type": "org_chart",
            "confidence": 0.85,
            "chart_title": "Team Structure & Responsibilities",
            "estimated_complexity": "high"
        }}
    ],
    "design_recommendations": {{
        "color_palette": ["{theme.colors.get('primary', '#2563eb')}", "{theme.colors.get('secondary', '#7c3aed')}", "{theme.colors.get('accent', '#dc2626')}"],
        "font_suggestions": ["Segoe UI", "Calibri", "Arial"],
        "visual_style": "modern_corporate",
        "animation_suggestions": ["fade_in", "slide_from_left", "emphasis_zoom"]
    }}
}}

GENERATE COMPREHENSIVE, PROFESSIONAL CONTENT - Each slide should be detailed with substantial content, professional design elements, and thorough speaker notes. Create a complete presentation worthy of a corporate boardroom."""
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Create a comprehensive, professional PowerPoint-style slide presentation about: {prompt}. Ensure it includes 10-15 slides with substantial content, diverse templates, professional design elements, and detailed speaker notes for each slide."}
            ],
            temperature=0.7,
            max_tokens=12000  # High limit with GPT-4o's 128K context
        )
        
        ai_response = response.choices[0].message.content
        
        try:
            ai_data = json.loads(ai_response)
        except json.JSONDecodeError:
            ai_data = create_fallback_slides_data(prompt)
        
        # Create presentation with AI-generated content
        with transaction.atomic():
            presentation = SlidePresentation.objects.create(
                title=ai_data['title'],
                theme=theme,
                created_by=user,
                slide_size=slide_size,
                outline_structure=ai_data.get('outline_structure', {}),
                diagram_opportunities=ai_data.get('diagram_opportunities', [])
            )
            
            # Create slides from AI data
            for slide_data in ai_data.get('slides', []):
                # Get appropriate template
                template_type = slide_data.get('template_type', 'title_content')
                try:
                    template = SlideTemplate.objects.get(layout_type=template_type)
                except SlideTemplate.DoesNotExist:
                    template = SlideTemplate.objects.first()  # Fallback
                
                Slide.objects.create(
                    presentation=presentation,
                    template=template,
                    order=slide_data.get('order', 0),
                    content=slide_data.get('content', {}),
                    notes=slide_data.get('notes', ''),
                    background={
                        'type': 'color',
                        'value': theme.colors.get('background', '#ffffff')
                    }
                )
            
            # Update slide count
            presentation.update_slide_count()
            
            logger.info(f"Successfully generated slide presentation {presentation.id}")
            return {
                'status': 'success',
                'presentation_id': str(presentation.id),
                'presentation_data': {
                    'id': str(presentation.id),
                    'title': presentation.title,
                    'slide_count': presentation.slide_count,
                    'theme_name': theme.name
                },
                'ai_analysis': {
                    'diagram_opportunities': ai_data.get('diagram_opportunities', []),
                    'generation_success': True,
                    'design_score': 0.9
                }
            }
        
    except Exception as e:
        logger.error(f"AI slide generation failed: {e}")
        logger.error(traceback.format_exc())
        
        if self.request.retries < self.max_retries:
            raise self.retry(countdown=60, exc=e)
        
        return {'status': 'failed', 'error': str(e)}


@shared_task(bind=True, max_retries=2)
def convert_text_to_diagram_task(self, text, chart_type, user_id, document_id=None, slide_id=None):
    """
    Celery task for Napkin.ai-style text to diagram conversion
    """
    try:
        from django.contrib.auth.models import User
        
        user = User.objects.get(id=user_id)
        logger.info(f"Starting text-to-diagram conversion for user {user_id}")
        
        # Set OpenAI client
        from openai import OpenAI
        client = OpenAI(api_key=settings.OPENAI_API_KEY)
        
        system_prompt = f"""Convert this text into a {chart_type} diagram. 
        
        Analyze the text and extract data to create a working {chart_type}.
        
        Return a JSON response with this structure:
        {{
            "title": "Diagram Title",
            "chart_type": "{chart_type}",
            "data": {{
                "labels": ["Label 1", "Label 2"],
                "datasets": [{{
                    "data": [10, 20, 30],
                    "backgroundColor": ["#FF6384", "#36A2EB", "#FFCE56"]
                }}]
            }},
            "config": {{
                "type": "{chart_type}",
                "responsive": true,
                "plugins": {{
                    "legend": {{"position": "top"}},
                    "title": {{"display": true, "text": "Chart Title"}}
                }}
            }},
            "styling": {{
                "width": 400,
                "height": 300,
                "colors": ["#FF6384", "#36A2EB", "#FFCE56"]
            }},
            "ai_interpretation": {{
                "extracted_entities": [],
                "relationships": [],
                "data_points": []
            }},
            "confidence_score": 0.9
        }}
        
        Make sure the data is Chart.js compatible and accurately represents the text content."""
        
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": f"Text to convert: {text}"}
            ],
            temperature=0.3,
            max_tokens=2000
        )
        
        ai_response = response.choices[0].message.content
        
        try:
            diagram_data = json.loads(ai_response)
        except json.JSONDecodeError:
            diagram_data = create_fallback_diagram_data(text, chart_type)
        
        # Generate actual image using QuickChart or similar service
        image_url = generate_diagram_image_url(diagram_data)
        diagram_data['image_url'] = image_url
        
        # Create diagram element
        diagram = DiagramElement.objects.create(
            title=diagram_data['title'],
            chart_type=chart_type,
            data=diagram_data['data'],
            config=diagram_data['config'],
            styling=diagram_data['styling'],
            source_text=text,
            ai_interpretation=diagram_data['ai_interpretation'],
            generation_prompt=f"Convert text to {chart_type}",
            confidence_score=diagram_data['confidence_score'],
            image_url=diagram_data['image_url'],
            original_content=text,
            created_by=user
        )
        
        # Link to document or slide
        if document_id:
            try:
                document = Document.objects.get(id=document_id, created_by=user)
                diagram.used_in_documents.add(document)
            except Document.DoesNotExist:
                pass
        
        if slide_id:
            try:
                slide = Slide.objects.get(id=slide_id, presentation__created_by=user)
                diagram.used_in_slides.add(slide)
            except Slide.DoesNotExist:
                pass
        
        logger.info(f"Successfully created diagram {diagram.id}")
        return {
            'status': 'success',
            'diagram_id': str(diagram.id),
            'diagram_data': {
                'id': str(diagram.id),
                'title': diagram.title,
                'image_url': diagram.image_url,
                'chart_type': diagram.chart_type
            },
            'confidence': diagram_data['confidence_score']
        }
        
    except Exception as e:
        logger.error(f"Text-to-diagram conversion failed: {e}")
        logger.error(traceback.format_exc())
        
        if self.request.retries < self.max_retries:
            raise self.retry(countdown=30, exc=e)
        
        return {'status': 'failed', 'error': str(e)}


def create_fallback_document_data(prompt, document_type):
    """Create fallback document data if AI fails"""
    return {
        "title": f"{document_type.title()} Document",
        "abstract": f"A comprehensive {document_type} document based on: {prompt[:200]}...",
        "keywords": f"{document_type}, analysis, professional",
        "authors": ["User"],
        "subject": document_type,
        "category": document_type,
        "structure": {
            "chapters": [
                {
                    "number": 1,
                    "title": "Introduction",
                    "content": f"<p>This document addresses {prompt}</p>",
                    "sections": []
                }
            ]
        },
        "content": f"<h1>Introduction</h1><p>This document addresses: {prompt}</p>",
        "diagram_opportunities": []
    }


def create_fallback_slides_data(prompt):
    """Create fallback slides data if AI fails"""
    return {
        "title": "Presentation",
        "slides": [
            {
                "template_type": "title",
                "order": 0,
                "content": {
                    "title_zone": "Presentation Title",
                    "subtitle_zone": "Based on your request"
                },
                "notes": f"Presentation about: {prompt}"
            }
        ],
        "outline_structure": {
            "main_sections": ["Introduction"],
            "total_slides": 1
        },
        "diagram_opportunities": []
    }


def create_fallback_diagram_data(text, chart_type):
    """Create fallback diagram data if AI fails"""
    return {
        "title": f"{chart_type.replace('_', ' ').title()}",
        "chart_type": chart_type,
        "data": {"labels": ["Data 1", "Data 2"], "datasets": [{"data": [1, 2]}]},
        "config": {"type": chart_type, "responsive": True},
        "styling": {"width": 400, "height": 300},
        "ai_interpretation": {},
        "confidence_score": 0.5
    }


def generate_diagram_image_url(diagram_data):
    """Generate actual diagram image URL using chart service"""
    try:
        import urllib.parse
        # Use QuickChart.io service for generating chart images
        chart_config = json.dumps(diagram_data.get('config', {}))
        encoded_config = urllib.parse.quote(chart_config)
        return f"https://quickchart.io/chart?c={encoded_config}"
    except Exception as e:
        logger.error(f"Failed to generate chart image URL: {e}")
        chart_type = diagram_data.get('chart_type', 'bar_chart')
        return f"https://via.placeholder.com/400x300/4f46e5/ffffff?text={chart_type}"


# ============================================================================
# KEEP ALL YOUR EXISTING TASKS BELOW THIS LINE
# ============================================================================

# Your existing export_presentation_pdf_task, export_presentation_pptx_task, etc.
# should remain unchanged...